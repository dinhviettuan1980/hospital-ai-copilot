"""
Weekly report management: upload a .pptx briefing deck, let Groq draft a
structured extraction against the report_period/metric_value/incident/
feedback/narrative_section schema, and let staff review/edit that draft
before it is written to Postgres. Also covers list/get/update/delete for
reports already in the DB.

The extraction step never touches the database - it only returns a draft
JSON blob. The client is responsible for sending back the (possibly edited)
draft to create_report/update_report.
"""
import io
import json
import re
import time

from pptx import Presentation

# This Groq org's on-demand tier caps openai/gpt-oss-120b at 8000 tokens/minute,
# which a single call covering both the metric table and the narrative text of a
# real weekly deck (~30-50 slides) can exceed. Splitting extraction into a
# table-focused call and a narrative-focused call keeps each request well under
# that budget instead of truncating (and silently dropping) report content.
MODEL = "openai/gpt-oss-120b"

FACILITY_CODES = ("CS1", "CS2", "TOTAL")

# Vietnamese text with lots of numbers/punctuation tokenizes at roughly
# 2.2-2.3 chars/token (measured against this org's actual Groq responses),
# not the ~4 chars/token rule of thumb for English. Char caps below are
# sized conservatively off that ratio so (system prompt + input +
# max_tokens reserved for the completion) stays under the 8000 TPM budget.
METRICS_MAX_COMPLETION_TOKENS = 2200
NARRATIVE_MAX_COMPLETION_TOKENS = 2400
MAX_TABLE_CHARS = 9000
MAX_NARRATIVE_CHARS = 9000


# Some decks put qualitative content (praise/complaint logs, "góp ý" tables) in
# an actual PowerPoint *table* rather than free text - if any of these appear
# in a table's cells, route the whole table to the narrative stream instead of
# the metrics stream, or it would silently be dropped (no metric_code matches
# a praise/complaint row, so a metrics-only prompt just ignores it).
_NARRATIVE_TABLE_KEYWORDS = (
    "góp ý", "khen", "sự cố", "khiếu nại", "phản ánh",
    "nguyên nhân và giải pháp", "được khen",
)


def extract_pptx_streams(file_bytes: bytes) -> tuple[str, str]:
    """Split a pptx into (table_text, narrative_text) - numeric tables carry
    the quantitative metrics, everything else (free text, and any table that
    looks qualitative) carries title/incidents/feedback/narrative - so each
    stream can be sent to Groq as a much smaller, focused prompt."""
    prs = Presentation(io.BytesIO(file_bytes))
    table_lines, narrative_lines = [], []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_table:
                rows_text = [
                    " | ".join(cell.text.strip() for cell in row.cells) for row in shape.table.rows
                ]
                joined = "\n".join(rows_text).lower()
                is_narrative = any(kw in joined for kw in _NARRATIVE_TABLE_KEYWORDS)
                target = narrative_lines if is_narrative else table_lines
                target.append(f"--- Slide {slide_idx} ---")
                target.extend(rows_text)
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    narrative_lines.append(f"--- Slide {slide_idx} ---")
                    narrative_lines.append(text)
    return "\n".join(table_lines), "\n".join(narrative_lines)


def load_metric_defs(conn):
    with conn.cursor() as cur:
        cur.execute("SELECT code, category, name, unit FROM metric ORDER BY category, code")
        return cur.fetchall()


def _build_metrics_prompt(metric_defs) -> str:
    metrics_txt = "\n".join(f"{m['code']}={m['name']}({m['unit']})" for m in metric_defs)
    return f"""Trích số liệu từ bảng báo cáo giao ban tuần bệnh viện. Trả về DUY NHẤT JSON, không giải thích:
{{"metrics": {{"<metric_code>": {{"CS1": number|null, "CS2": number|null, "TOTAL": number|null, "note": string|null}}}}}}

metric_code hợp lệ (bỏ qua số liệu không khớp code nào):
{metrics_txt}

Quy tắc quan trọng:
- Dòng tổng (vd "...2 cơ sở") + 2 dòng con "Cơ sở 1"/"Cơ sở 2" ngay dưới -> dòng tổng=TOTAL, 2 dòng con=CS1/CS2.
- Bảng có nhiều cột số liệu (TB tháng trước, tuần báo cáo, chênh lệch/%) -> CHỈ lấy cột khớp tên/khoảng ngày
  tuần báo cáo, bỏ qua cột tháng trước và cột chênh lệch/%. Cột "tuần báo cáo" là cột có tiêu đề dạng
  "Tuần DD/M->DD/M/YYYY" hoặc tương tự - đây LUÔN LÀ CỘT THỨ 2 trong nhóm 3 cột (tháng trước | tuần này |
  chênh lệch), KHÔNG PHẢI cột đầu tiên.
- VÍ DỤ CỤ THỂ: dòng dữ liệu "I. | Tổng số BN khám bệnh 2 cơ sở | 12.288 | 11.489 | 799" với tiêu đề cột
  "T7/2026 | Tuần 19/8->25/8/2026 | Chênh lệch" -> giá trị ĐÚNG cần lấy là 11.489 (cột thứ 2, khớp tuần báo
  cáo), TUYỆT ĐỐI KHÔNG lấy 12.288 (cột thứ 1, chỉ là số liệu trung bình tháng trước, không phải số của tuần
  báo cáo này).
- "12.288" = 12288 (dấu chấm là hàng nghìn kiểu VN, không phải thập phân).
- Không chắc chắn -> để null, không bịa số (dữ liệu dùng cho Ban Giám đốc bệnh viện)."""


_NARRATIVE_PROMPT = """Bạn là bot trích xuất dữ liệu định tính từ báo cáo giao ban tuần của bệnh viện.
Người dùng sẽ đưa cho bạn nội dung text (tiêu đề, sự cố, khiếu nại/khen ngợi, tường thuật theo khoa/phòng)
đã trích thô từ file PowerPoint. Đọc và trả về DUY NHẤT một JSON object theo đúng schema sau, không thêm
giải thích:

{
  "label": "Tuần DD/MM - DD/MM/YYYY",
  "start_date": "YYYY-MM-DD",
  "end_date": "YYYY-MM-DD",
  "incidents": [
    {"department": string, "incident_date": string, "description": string, "cause": string,
      "corrective_action": string, "resolved": boolean, "severity": "low"|"medium"|"high"}
  ],
  "feedback": [
    {"date": string, "department": string, "type": "complaint"|"praise", "content": string,
      "cause": string|null, "resolution": string|null}
  ],
  "narrative_sections": [
    {"section_name": string, "content": string}
  ]
}

label/start_date/end_date thường lấy từ slide tiêu đề đầu báo cáo (dạng "Từ ngày DD/MM/YYYY đến DD/MM/YYYY").
Nếu không chắc chắn hoặc không tìm thấy thông tin nào, bỏ qua thay vì bịa ra - đây là dữ liệu dùng cho
Ban Giám đốc bệnh viện, sai lệch là không chấp nhận được."""


def _chat_json(groq_client, system_prompt: str, user_content: str, max_tokens: int, retries: int = 1) -> dict:
    for attempt in range(retries + 1):
        try:
            response = groq_client.chat.completions.create(
                model=MODEL,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_content},
                ],
                response_format={"type": "json_object"},
                max_tokens=max_tokens,
                temperature=0.1,
                reasoning_effort="low",
            )
            content = response.choices[0].message.content
            if not content:
                raise ValueError("Groq trả về nội dung rỗng (có thể do hết token cho phần suy luận nội bộ)")
            return json.loads(content)
        except Exception as e:
            msg = str(e)
            # 429 = temporarily out of budget in the current rate-limit window, worth
            # a short wait and retry. 413 means this single request's own size (prompt
            # + max_tokens) permanently exceeds the tier's per-request cap - retrying
            # the identical payload will never succeed, so fail fast instead.
            if attempt < retries and "429" in msg and "413" not in msg:
                time.sleep(20)
                continue
            raise


# Hard ceiling on chunk count per stream, so a pathologically large deck fails
# loudly (flagged as truncated) instead of firing dozens of sequential Groq
# calls and taking many minutes.
MAX_CHUNKS = 6


def _chunk_by_slide(text: str, max_chars: int) -> tuple[list[str], bool]:
    """Split text into <= max_chars pieces on slide boundaries (never mid-row),
    so a large deck can be extracted across several smaller Groq calls instead
    of truncating (and losing) the tail. Returns (chunks, was_capped)."""
    if len(text) <= max_chars:
        return ([text] if text.strip() else []), False
    blocks = re.split(r"(?=--- Slide \d+ ---)", text)
    chunks: list[str] = []
    current = ""
    for block in blocks:
        if not block.strip():
            continue
        if current and len(current) + len(block) > max_chars:
            chunks.append(current)
            current = block
        else:
            current += block
    if current:
        chunks.append(current)
    if len(chunks) > MAX_CHUNKS:
        return chunks[:MAX_CHUNKS], True
    return chunks, False


def _merge_metrics(base: dict, addition: dict) -> dict:
    for code, values in addition.items():
        existing = base.get(code)
        if existing is None or all(existing.get(f) is None for f in FACILITY_CODES):
            base[code] = values
    return base


def extract_report_data(table_text: str, narrative_text: str, groq_client, conn) -> dict:
    metric_defs = load_metric_defs(conn)
    metrics_prompt = _build_metrics_prompt(metric_defs)

    table_chunks, table_capped = _chunk_by_slide(table_text, MAX_TABLE_CHARS)
    metrics: dict = {}
    for i, chunk in enumerate(table_chunks):
        chunk_result = _chat_json(groq_client, metrics_prompt, chunk, METRICS_MAX_COMPLETION_TOKENS)
        _merge_metrics(metrics, chunk_result.get("metrics", {}))
        if i < len(table_chunks) - 1:
            time.sleep(8)  # let the per-minute token bucket refill before the next chunk

    narrative_chunks, narrative_capped = _chunk_by_slide(narrative_text, MAX_NARRATIVE_CHARS)
    label = start_date = end_date = ""
    incidents, feedback, narrative_sections = [], [], []
    for i, chunk in enumerate(narrative_chunks):
        chunk_result = _chat_json(groq_client, _NARRATIVE_PROMPT, chunk, NARRATIVE_MAX_COMPLETION_TOKENS)
        label = label or chunk_result.get("label") or ""
        start_date = start_date or chunk_result.get("start_date") or ""
        end_date = end_date or chunk_result.get("end_date") or ""
        incidents.extend(chunk_result.get("incidents", []))
        feedback.extend(chunk_result.get("feedback", []))
        narrative_sections.extend(chunk_result.get("narrative_sections", []))
        if i < len(narrative_chunks) - 1:
            time.sleep(8)

    return {
        "label": label,
        "start_date": start_date,
        "end_date": end_date,
        "metrics": metrics,
        "incidents": incidents,
        "feedback": feedback,
        "narrative_sections": narrative_sections,
        "truncated": {"table": table_capped, "narrative": narrative_capped},
    }


def _facility_ids(conn) -> dict:
    with conn.cursor() as cur:
        cur.execute("SELECT id, code FROM facility")
        return {row["code"]: row["id"] for row in cur.fetchall()}


def _metric_ids(conn) -> dict:
    with conn.cursor() as cur:
        cur.execute("SELECT id, code FROM metric")
        return {row["code"]: row["id"] for row in cur.fetchall()}


def list_reports(conn):
    with conn.cursor() as cur:
        cur.execute(
            """
            SELECT rp.id, rp.label, rp.start_date, rp.end_date, rp.source_file,
                   (SELECT count(*) FROM incident i WHERE i.report_period_id = rp.id) AS incident_count,
                   (SELECT count(*) FROM feedback f WHERE f.report_period_id = rp.id) AS feedback_count
            FROM report_period rp
            ORDER BY rp.start_date DESC
            """
        )
        return cur.fetchall()


def find_overlapping_reports(conn, start_date: str, end_date: str, exclude_id: int | None = None):
    with conn.cursor() as cur:
        cur.execute(
            """
            SELECT id, label, start_date, end_date, source_file FROM report_period
            WHERE start_date <= %s AND end_date >= %s AND (%s::int IS NULL OR id != %s)
            ORDER BY start_date
            """,
            (end_date, start_date, exclude_id, exclude_id),
        )
        rows = cur.fetchall()
        for row in rows:
            row["start_date"] = str(row["start_date"])
            row["end_date"] = str(row["end_date"])
        return rows


def get_report(conn, report_id: int):
    with conn.cursor() as cur:
        cur.execute("SELECT * FROM report_period WHERE id = %s", (report_id,))
        rp = cur.fetchone()
        if rp is None:
            return None

        cur.execute(
            """
            SELECT m.code, f.code AS facility_code, mv.value, mv.note
            FROM metric_value mv
            JOIN metric m ON m.id = mv.metric_id
            JOIN facility f ON f.id = mv.facility_id
            WHERE mv.report_period_id = %s
            """,
            (report_id,),
        )
        metrics: dict = {}
        for row in cur.fetchall():
            entry = metrics.setdefault(row["code"], {"CS1": None, "CS2": None, "TOTAL": None, "note": None})
            entry[row["facility_code"]] = row["value"]
            if row["note"]:
                entry["note"] = row["note"]

        cur.execute(
            """
            SELECT department, incident_date, description, cause, corrective_action, resolved, severity
            FROM incident WHERE report_period_id = %s ORDER BY id
            """,
            (report_id,),
        )
        incidents = cur.fetchall()

        cur.execute(
            """
            SELECT date, department, type, content, cause, resolution
            FROM feedback WHERE report_period_id = %s ORDER BY id
            """,
            (report_id,),
        )
        feedback = cur.fetchall()

        cur.execute(
            "SELECT section_name, content FROM narrative_section WHERE report_period_id = %s ORDER BY id",
            (report_id,),
        )
        narrative_sections = cur.fetchall()

    return {
        "id": rp["id"],
        "label": rp["label"],
        "start_date": str(rp["start_date"]),
        "end_date": str(rp["end_date"]),
        "source_file": rp["source_file"],
        "metrics": metrics,
        "incidents": incidents,
        "feedback": feedback,
        "narrative_sections": narrative_sections,
    }


def _insert_children(conn, report_id: int, data: dict, facility_ids: dict, metric_ids: dict):
    with conn.cursor() as cur:
        for code, values in data.get("metrics", {}).items():
            metric_id = metric_ids.get(code)
            if metric_id is None:
                continue
            note = values.get("note")
            for fac_code in FACILITY_CODES:
                value = values.get(fac_code)
                if value is None:
                    continue
                cur.execute(
                    "INSERT INTO metric_value (report_period_id, facility_id, metric_id, value, note) "
                    "VALUES (%s, %s, %s, %s, %s)",
                    (report_id, facility_ids[fac_code], metric_id, value, note),
                )

        for inc in data.get("incidents", []):
            cur.execute(
                "INSERT INTO incident (report_period_id, department, incident_date, description, cause, "
                "corrective_action, resolved, severity) VALUES (%s, %s, %s, %s, %s, %s, %s, %s)",
                (
                    report_id, inc.get("department"), inc.get("incident_date"), inc.get("description"),
                    inc.get("cause"), inc.get("corrective_action"), bool(inc.get("resolved")), inc.get("severity"),
                ),
            )

        for fb in data.get("feedback", []):
            cur.execute(
                "INSERT INTO feedback (report_period_id, date, department, type, content, cause, resolution) "
                "VALUES (%s, %s, %s, %s, %s, %s, %s)",
                (
                    report_id, fb.get("date"), fb.get("department"), fb.get("type"), fb.get("content"),
                    fb.get("cause"), fb.get("resolution"),
                ),
            )

        for sec in data.get("narrative_sections", []):
            cur.execute(
                "INSERT INTO narrative_section (report_period_id, section_name, content) VALUES (%s, %s, %s)",
                (report_id, sec.get("section_name"), sec.get("content")),
            )


def create_report(conn, data: dict, source_file: str) -> int:
    facility_ids = _facility_ids(conn)
    metric_ids = _metric_ids(conn)
    try:
        with conn.cursor() as cur:
            cur.execute(
                "INSERT INTO report_period (label, start_date, end_date, source_file) "
                "VALUES (%s, %s, %s, %s) RETURNING id",
                (data["label"], data["start_date"], data["end_date"], source_file),
            )
            report_id = cur.fetchone()["id"]
        _insert_children(conn, report_id, data, facility_ids, metric_ids)
        conn.commit()
        return report_id
    except Exception:
        conn.rollback()
        raise


def _delete_children(conn, report_id: int):
    with conn.cursor() as cur:
        cur.execute("DELETE FROM metric_value WHERE report_period_id = %s", (report_id,))
        cur.execute("DELETE FROM incident WHERE report_period_id = %s", (report_id,))
        cur.execute("DELETE FROM feedback WHERE report_period_id = %s", (report_id,))
        cur.execute("DELETE FROM narrative_section WHERE report_period_id = %s", (report_id,))


def update_report(conn, report_id: int, data: dict) -> bool:
    facility_ids = _facility_ids(conn)
    metric_ids = _metric_ids(conn)
    try:
        with conn.cursor() as cur:
            cur.execute(
                "UPDATE report_period SET label = %s, start_date = %s, end_date = %s WHERE id = %s",
                (data["label"], data["start_date"], data["end_date"], report_id),
            )
            if cur.rowcount == 0:
                conn.rollback()
                return False
        _delete_children(conn, report_id)
        _insert_children(conn, report_id, data, facility_ids, metric_ids)
        conn.commit()
        return True
    except Exception:
        conn.rollback()
        raise


def delete_report(conn, report_id: int) -> bool:
    try:
        _delete_children(conn, report_id)
        with conn.cursor() as cur:
            cur.execute("DELETE FROM report_period WHERE id = %s", (report_id,))
            deleted = cur.rowcount > 0
        conn.commit()
        return deleted
    except Exception:
        conn.rollback()
        raise
